import pandas as pd
from pptx import Presentation
from utils import ConfigController, find_layout_by_name, read_csvs_to_dfs
from chart_handler import ChartHandler
from pptx.util import Inches


class deckGen:
    def __init__(self, mode, config, df=None):
        self.config = ConfigController(config)
        self.Presentation = self.read_template(self.config)
        self.slides = self.config.get_slides()

        match mode:
            case "test":     
                self.config.display_config()
                from data_handlers import TestDataGenerator
                self.dataGenerator = TestDataGenerator(num_cats=2, num_measures=3, num_personas=10, include_timeseries=True)
        
            case "csv":
                from data_handlers import CsvDataGenerator
                self.dataGenerator = CsvDataGenerator(self.config.get_data())
            
            case "dataframe":
                print("Processing DataFrame.. One Moment..")
                if config:
                    self.dataGenerator = df
                else: 
                    raise FileNotFoundError('File Not Found: {} '.format(config.get_template_path()))
            case _:
                print("The mode input is not currently supported. Please check the input mode.")

    def generate_deck(self):
        # this block might be very sensitive, it handles the title slide 
        title_slide =  self.Presentation.slides[0] #this seems  like the only reliable way to get the title slide
        title = title_slide.shapes.title
        title.text = self.config.get_presentation_title()
        sub_title = title_slide.placeholders[1]
        sub_title.text = self.config.get_presentation_author()

        # this block is iterating through the slides listed in the  yaml config
        for idx, slide_ref in enumerate(self.slides):
            slide = Slide(slide_ref, idx)
            slide.fill_slide(self.Presentation, self.dataGenerator)

        self.Presentation.save(self.config.get_output_path())
        return

    def read_template(self, config):
        try:        
            presentation = Presentation(config.get_template_path())
            return presentation
        except:
            raise FileNotFoundError('File Not Found: {} '.format(config.get_template_path()))

class Slide:
    def __init__(self, slide_config, idx, data=None):
        self.title = slide_config['Title']
        self.configLayout = slide_config["Layout"]
        self.charts = slide_config['Charts']
        self.slide_idx = idx

    def fill_slide(self, prs, dataGenerator):
        slide_layout = find_layout_by_name(prs, self.configLayout)
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = self.title

        for _, chartConfig  in enumerate(self.charts):
            print("{} Placeholders Available for Slide {}".format(len(slide.placeholders), self.slide_idx+1 ))
            chartController = ChartHandler(self.slide_idx, chartConfig, dataGenerator)

            ###
            #pptx and the python  api have a quirk with the identification of ids
            # when in the selection pane in pptx, the id of the placeholder is in
            # reverse order of the indexs of the pane, to reference the first place holder
            # we caluclate the idx below.  
            ###
            idx = len(slide.placeholders) - chartConfig["Location"] +2
            for _, placeholder in enumerate(slide.placeholders):
                if idx == placeholder.shape_id:
                    
                    if chartConfig.get('Args', None).get('chartType', None) == 'excel':
                        chartController.getExcel(slide, placeholder)      
                    else:
                        chartController.insertChart(placeholder)

        return 
    
    
        


if __name__ == "__main__":
    # print('Handling Test Config')
    # dG = deckGen("test", "test/TestConfig")
    # #TODO: handle aargs to route test mode, db, file, etc
    # dG.generate_deck()

    print('Handling Csv Config')
    dG = deckGen("csv", "test/CsvConfig")
    #TODO: handle aargs to route test mode, db, file, etc
    dG.generate_deck()